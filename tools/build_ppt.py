# -*- coding: utf-8 -*-
"""
Build a clean portfolio PPTX that mirrors the joannekim0420.github.io theme:
white canvas, charcoal text, per-section accent colors, Montserrat (Latin)
+ Malgun Gothic (Korean). Run: python tools/build_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# ---- theme ----------------------------------------------------------------
FG       = (26, 27, 31)      # #1a1b1f
MUTED    = (67, 70, 77)      # #43464d
SOFTGRAY = (117, 134, 150)   # softened
BG_SOFT  = (245, 247, 250)   # #f5f7fa
BORDER   = (228, 235, 243)   # #e4ebf3
WHITE    = (255, 255, 255)
BRAND    = (31, 41, 55)      # #1f2937 charcoal
HDC      = (200, 16, 46)     # #c8102e
UNIST    = (0, 102, 179)     # #0066b3
NCC      = (79, 70, 229)     # #4f46e5
SMU      = (71, 85, 105)     # #475569 neutral

LATIN = "Montserrat"
EA    = "Malgun Gothic"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _ea(run, ea=EA):
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    el = rPr.find(qn('a:ea'))
    if el is None:
        el = rPr.makeelement(qn('a:ea'), {})
        (latin.addnext(el) if latin is not None else rPr.append(el))
    el.set('typeface', ea)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(*fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor(*line); sp.line.width = Pt(line_w)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.06, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,bold,color) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = RGBColor(*color); r.font.name = LATIN
            _ea(r)
    return tb


def img_fit(s, path, x, y, max_w, max_h, align="center", valign="middle"):
    p = os.path.join(ROOT, path)
    if not os.path.exists(p):
        return None
    iw, ih = Image.open(p).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    if align == "center":
        x = x + (max_w - w) // 2
    elif align == "right":
        x = x + (max_w - w)
    if valign == "middle":
        y = y + (max_h - h) // 2
    elif valign == "bottom":
        y = y + (max_h - h)
    return s.shapes.add_picture(p, x, y, w, h)


PAGE = 0
TOTAL = 13


def chrome(s, eyebrow, title, accent, sub=None):
    """Standard content-slide header + footer."""
    global PAGE
    PAGE += 1
    rect(s, 0, 0, Inches(0.16), EMU_H, fill=accent)            # left accent spine
    text(s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.3),
         [[(eyebrow, 12, True, accent)]])
    text(s, Inches(0.7), Inches(0.82), Inches(11.9), Inches(0.7),
         [[(title, 30, True, FG)]])
    rect(s, Inches(0.72), Inches(1.52), Inches(0.7), Pt(3), fill=accent)
    if sub:
        text(s, Inches(0.7), Inches(1.66), Inches(11.9), Inches(0.4),
             [[(sub, 13, False, MUTED)]])
    # footer
    text(s, Inches(0.7), Inches(7.06), Inches(7), Inches(0.3),
         [[("김종은 Jongeun Kim · AI Engineer", 9, False, SOFTGRAY)]])
    text(s, Inches(11.6), Inches(7.06), Inches(1.0), Inches(0.3),
         [[(f"{PAGE:02d} / {TOTAL:02d}", 9, False, SOFTGRAY)]], align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, fill=BG_SOFT, line=BORDER):
    return rect(s, x, y, w, h, fill=fill, line=line, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def chip(s, x, y, label, accent):
    w = Inches(0.10 + 0.085 * len(label))
    c = rect(s, x, y, w, Inches(0.30), fill=BG_SOFT, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(10.5); r.font.bold = True
    r.font.color.rgb = RGBColor(*accent); r.font.name = LATIN; _ea(r)
    return x + w + Inches(0.12)


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=WHITE)
rect(s, 0, 0, Inches(0.22), EMU_H, fill=BRAND)
# right photo panel
img_fit(s, "assets/img/profile/profile.jpg", Inches(9.0), Inches(1.7),
        Inches(3.4), Inches(4.1))
text(s, Inches(0.9), Inches(1.85), Inches(8.0), Inches(0.4),
     [[("PORTFOLIO", 13, True, HDC)]])
text(s, Inches(0.9), Inches(2.25), Inches(8.0), Inches(1.2),
     [[("김종은 ", 46, True, FG), ("Jongeun Kim", 46, True, FG)]])
text(s, Inches(0.9), Inches(3.25), Inches(8.0), Inches(0.5),
     [[("AI Engineer · Multimodal AI · VLM · LLM", 18, True, MUTED)]])
text(s, Inches(0.92), Inches(3.95), Inches(7.6), Inches(0.9),
     [[("VLM·LLM·비디오 언더스탠딩 — 연구부터 현장 적용까지 멀티모달 AI 개발", 14, False, MUTED)],
      [("UNIST 인공지능대학원 석사 · 논문 7편 · 특허 7건", 14, False, MUTED)]],
     line_spacing=1.2, space_after=6)
rect(s, Inches(0.92), Inches(5.0), Inches(0.7), Pt(3), fill=HDC)
text(s, Inches(0.92), Inches(5.2), Inches(7.6), Inches(1.0),
     [[("GitHub  github.com/joannekim0420", 12, False, MUTED)],
      [("Scholar  PHEuuWoAAAAJ   ·   Email  JongeunKim@r114.com", 12, False, MUTED)]],
     line_spacing=1.25, space_after=4)

# ===========================================================================
# 2. ABOUT / OVERVIEW
# ===========================================================================
s = slide()
chrome(s, "ABOUT", "소개", BRAND)
text(s, Inches(0.7), Inches(2.1), Inches(7.0), Inches(2.2),
     [[("홈 AI 에이전트, 엘리베이터 상황 분석, 화재 탐지 등 실환경", 15, False, MUTED)],
      [("멀티모달 AI 과제를 PoC부터 현장 적용까지 수행", 15, False, MUTED)],
      [("", 6, False, MUTED)],
      [("UNIST 인공지능대학원 석사 (멀티모달·NLP)", 15, False, MUTED)],
      [("연구 성과를 논문 7편과 특허 7건으로 권리화", 15, False, MUTED)]],
     line_spacing=1.3, space_after=6)
# stat cards
stats = [("7", "논문 Publications", UNIST), ("7", "특허 Patents", HDC),
         ("M.S.", "UNIST AI 석사", NCC), ("4+", "현장 프로젝트", SMU)]
cx = Inches(8.0)
cw, gap = Inches(2.45), Inches(0.18)
for i, (big, lab, col) in enumerate(stats):
    yy = Inches(2.0) + i * Inches(1.12)
    card(s, cx, yy, cw, Inches(0.95))
    rect(s, cx, yy + Inches(0.18), Inches(0.10), Inches(0.6), fill=col)
    text(s, cx + Inches(0.28), yy + Inches(0.1), Inches(2.1), Inches(0.5),
         [[(big, 24, True, col)]])
    text(s, cx + Inches(0.28), yy + Inches(0.58), Inches(2.1), Inches(0.3),
         [[(lab, 11, True, MUTED)]])

# ===========================================================================
# 3. CAREER
# ===========================================================================
s = slide()
chrome(s, "CAREER", "경력", BRAND)
careers = [
    ("assets/img/CI/Hdclabs.png", "HDC LABS", "R&D 센터 AI Lab", "2024.05 ~ 현재", HDC),
    ("assets/img/CI/unist.png", "UNIST", "인공지능 석사", "2022.08 ~ 2024.08", UNIST),
    ("assets/img/CI/NC.jpeg", "NC SOFT", "기계번역데이터팀 인턴", "2021.07 ~ 2022.01", NCC),
    ("assets/img/CI/sangmyung.jpg", "상명대학교", "휴먼지능정보공학과 학사", "2017.02 ~ 2022.02", SMU),
]
cw = Inches(2.85); gap = Inches(0.18); x0 = Inches(0.7); y0 = Inches(2.2); ch = Inches(3.7)
for i, (logo, name, role, period, col) in enumerate(careers):
    x = x0 + i * (cw + gap)
    card(s, x, y0, cw, ch)
    rect(s, x, y0, cw, Inches(0.14), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    img_fit(s, logo, x + Inches(0.35), y0 + Inches(0.5), cw - Inches(0.7), Inches(0.95))
    text(s, x + Inches(0.2), y0 + Inches(1.7), cw - Inches(0.4), Inches(0.5),
         [[(name, 17, True, FG)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.2), y0 + Inches(2.2), cw - Inches(0.4), Inches(0.8),
         [[(role, 12, False, MUTED)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.2), y0 + Inches(3.15), cw - Inches(0.4), Inches(0.3),
         [[(period, 11, True, col)]], align=PP_ALIGN.CENTER)

# ===========================================================================
# 4. HDC overview
# ===========================================================================
s = slide()
chrome(s, "HDC LABS · 2024.05 ~", "VLM·LLM 기반 멀티모달 AI 현장 적용", HDC,
       sub="R&D 센터 AI Lab · AI Engineer — 특허 6건 등록·1건 출원, 국제 학회 논문 4편")
projects = [
    ("차세대 홈 AI Agent 시스템", "2025.06 ~", "기여도 35%",
     "VideoLLM+LLM 거주자 행동 인식·개인화, 정확도 65%→81%, Jetson NX 4bit 엣지 배포"),
    ("엘리베이터 내부 상황 분석 PoC", "2025.02 – 2025.03", "기여도 30%",
     "VLM 파인튜닝 반려견 탑승 인식 57%→87.5%, 특허 등록, 현대엘리베이터 협업"),
    ("지하주차장 화재 탐지", "2024.09 – 2025.04", "기여도 40%",
     "YOLO+VLM(Florence2) 통합, 탐지 6.51s→3.03s·Recall 0.835→0.904, ICLR 2025 워크샵"),
]
y = Inches(2.3)
for name, period, contrib, desc in projects:
    card(s, Inches(0.7), y, Inches(11.9), Inches(1.35))
    rect(s, Inches(0.7), y, Inches(0.12), Inches(1.35), fill=HDC)
    text(s, Inches(1.05), y + Inches(0.16), Inches(8.7), Inches(0.4),
         [[(name, 16, True, FG)]])
    text(s, Inches(1.05), y + Inches(0.62), Inches(10.8), Inches(0.6),
         [[(desc, 12.5, False, MUTED)]], line_spacing=1.15)
    text(s, Inches(9.9), y + Inches(0.16), Inches(2.55), Inches(0.35),
         [[(period, 11, True, HDC)]], align=PP_ALIGN.RIGHT)
    text(s, Inches(9.9), y + Inches(0.5), Inches(2.55), Inches(0.3),
         [[(contrib, 11, True, MUTED)]], align=PP_ALIGN.RIGHT)
    y += Inches(1.5)


def project_slide(eyebrow, title, accent, sub, spec, phases, image):
    s = slide()
    chrome(s, eyebrow, title, accent, sub=sub)
    # left: spec + phases
    y = Inches(2.25)
    for k, v in spec:
        text(s, Inches(0.7), y, Inches(1.3), Inches(0.4), [[(k, 11.5, True, accent)]])
        text(s, Inches(2.0), y, Inches(5.4), Inches(0.5), [[(v, 11.5, False, MUTED)]],
             line_spacing=1.1)
        y += Inches(0.52)
    y += Inches(0.1)
    for head, items in phases:
        text(s, Inches(0.7), y, Inches(6.8), Inches(0.3), [[(head, 13, True, FG)]])
        y += Inches(0.36)
        for it in items:
            text(s, Inches(0.92), y, Inches(6.5), Inches(0.5),
                 [[("•  ", 11.5, True, accent), (it, 11.5, False, MUTED)]],
                 line_spacing=1.12)
            y += Inches(0.46)
        y += Inches(0.08)
    # right: image card
    card(s, Inches(7.9), Inches(2.25), Inches(4.7), Inches(4.4), fill=WHITE)
    img_fit(s, image, Inches(8.05), Inches(2.4), Inches(4.4), Inches(4.1))
    return s


# 5. Home AI Agent
project_slide(
    "HDC LABS · PROJECT 01", "차세대 홈 AI Agent 시스템", HDC,
    "2025.06 ~ 진행 중 · 기여도 35% (프로젝트 매니징·데이터·모델링·파이프라인)",
    [("적용 환경", "실제 주거 공간 · 모델하우스"),
     ("연동 대상", "카메라·IoT 가전·월패드 (자체 플랫폼)"),
     ("핵심 기능", "VideoLLM 행동·물체·얼굴 인식 + LLM 에이전트 개인화")],
    [("주요 성과", [
        "단일 라벨 출력 제약·평가 기준 프롬프트 개선으로 행동 분류 정확도 65% → 81%",
        "Jetson NX 4bit/NF4 탑재, 5초 클립 평균 4.2초 처리로 준실시간 적용 가능성 확인",
        "촬영→라벨링→검증→학습→배포 5단계 자동화 툴 개발",
        "행동 인식 기반 IoT 가전 제어 특허 확보 (KR 10-2971047)"])],
    "assets/img/ext/hdc_aihome2.png")

# 6. Elevator
project_slide(
    "HDC LABS · PROJECT 02", "엘리베이터 내부 상황 분석 PoC", HDC,
    "2025.02 – 2025.03 · 기여도 30% (반려견 탑승 인식 모델·알고리즘)",
    [("적용 환경", "엘리베이터 카 내부 상단 모서리 CCTV 시점"),
     ("핵심 기능", "VLM 파인튜닝 반려견 탑승 인식 + 목줄 끼임 예방"),
     ("협업", "현대엘리베이터")],
    [("주요 성과", [
        "반려견 탑승 탐지 정확도 57% → 87.5%",
        "반려견·유사 객체(인형 등) 구분 정확도 91%",
        "엘리베이터 내 반려동물 탐지 특허 등록 (KR 10-2876730)",
        "HDC랩스·HDC현산·현대엘리베이터 업무협약, WACVW 2026 후속 논문"])],
    "assets/img/ext/hdc_elevator.png")

# 7. Fire detection
project_slide(
    "HDC LABS · PROJECT 03", "지하주차장 화재 탐지", HDC,
    "2024.09 – 2025.04 · 기여도 40% (VLM 통합 알고리즘·모델 선정)",
    [("적용 환경", "지하주차장 (전기차 충전구역)"),
     ("핵심 기능", "YOLO 기반 연기/불 탐지 + VLM 기반 상황 인식"),
     ("협업", "수호이미지랩스 (CCTV 전문기업)")],
    [("주요 성과", [
        "최초 화재 탐지 시간 6.51s → 3.03s, Recall 0.835 → 0.904",
        "YOLO+VLM(Florence2) 통합 상황 인식 알고리즘 제안",
        "ICLR 2025 ICBINB 워크샵 논문 게재",
        "복합모델 주차장 화재 탐지 특허 등록 (KR 10-2857163)",
        "광명 센트럴 아이파크 AI CCTV 상용 도입"])],
    "assets/img/ext/hdc_firedetection.png")

# ===========================================================================
# 8. UNIST
# ===========================================================================
s = slide()
chrome(s, "UNIST · 2022.08 – 2024.08", "인공지능대학원 석사 — NLP·멀티모달", UNIST,
       sub="졸업 논문: Efficient Multilingual Multimodal Fusion via Contrastive Learning")
uprojects = [
    ("고객 발화 Variation 텍스트 생성 (LG)", "생성형 AI 데이터 증강 · AI Judge 정량 평가 95.77% · LG H&A 산학과제 연장"),
    ("다국어 멀티모달 융합 (석사 졸업논문)", "기계번역 14개 언어 증강 · CLIP+XLM-RoBERTa 대조학습 · KSBE 2025"),
    ("슬로건 자동 생성 (심도컴퍼니)", "노이즈 섭동 기반 다양성 제어 생성 · 기존 대비 SOTA · CIKM 2023 1저자"),
    ("CVPR'22 LOVEU 챌린지 (Pyler)", "영상·스크립트·질의 멀티모달 QA · 평균 Recall +0.075 · 챌린지 2위"),
]
y = Inches(2.45)
for name, desc in uprojects:
    card(s, Inches(0.7), y, Inches(11.9), Inches(1.0))
    rect(s, Inches(0.7), y, Inches(0.12), Inches(1.0), fill=UNIST)
    text(s, Inches(1.05), y + Inches(0.14), Inches(11.2), Inches(0.4),
         [[(name, 15, True, FG)]])
    text(s, Inches(1.05), y + Inches(0.55), Inches(11.2), Inches(0.4),
         [[(desc, 12, False, MUTED)]])
    y += Inches(1.12)

# ===========================================================================
# 9. NC SOFT
# ===========================================================================
s = slide()
chrome(s, "NC SOFT · 2021.07 – 2022.01", "다국어 기계 번역 문장 품질 평가", NCC,
       sub="기계번역데이터팀 인턴 · 기여도 100%")
text(s, Inches(0.7), Inches(2.3), Inches(6.7), Inches(0.9),
     [[("원문-번역문 간 의미적 일치도를 수치화하여 사람 평가를 대체·보조하는", 13.5, False, MUTED)],
      [("자동 평가 방법 연구", 13.5, False, MUTED)]], line_spacing=1.25, space_after=4)
for i, it in enumerate([
        "사전 학습 다국어 문장 인코더로 원문-번역문 쌍을 같은 임베딩 공간에 정렬",
        "Cosine Similarity 기반 파인튜닝으로 사람 평가 점수와 상관도 향상",
        "기계번역 데이터팀 내부 기술 발표로 공유"]):
    yy = Inches(3.35) + i * Inches(0.6)
    text(s, Inches(0.92), yy, Inches(6.5), Inches(0.55),
         [[("•  ", 12, True, NCC), (it, 12, False, MUTED)]], line_spacing=1.12)
card(s, Inches(7.9), Inches(2.25), Inches(4.7), Inches(4.4), fill=WHITE)
img_fit(s, "assets/img/ext/nc_quality.png", Inches(8.05), Inches(2.4), Inches(4.4), Inches(4.1))

# ===========================================================================
# 10. PUBLICATIONS
# ===========================================================================
s = slide()
chrome(s, "PUBLICATIONS", "논문", BRAND, sub="국제 학회 6편 · 국내 학회 1편")
pubs = [
    ("WACVW 2026", "From Prompts to Deployment: Auto-Curated Domain-Specific Dataset Generation via Diffusion Models", "Oral"),
    ("ICCVW 2025", "Your Super Resolution Model is not Enough for Tackling Real-World Scenarios", ""),
    ("ICLRW 2025", "An Integrated YOLO and VLM System for Fire Detection in Enclosed Environments", ""),
    ("SIGGRAPH Asia 2024", "OverallNet: Scale-Arbitrary Lightweight SR Model for 360° Panoramic Images", "Poster"),
    ("CIKM 2023", "Effective Slogan Generation with Noise Perturbation", "1저자"),
    ("CVPRW 2022", "Technical Report for CVPR 2022 LOVEU AQTC Challenge", "2nd Place"),
    ("KSBE 2025", "Contrastive Learning for Efficient Multilingual Multimodal Fusion", "국내"),
]
y = Inches(2.15)
for venue, title, tag in pubs:
    rect(s, Inches(0.7), y + Inches(0.04), Inches(2.0), Inches(0.46),
         fill=BG_SOFT, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.7), y + Inches(0.04), Inches(2.0), Inches(0.46),
         [[(venue, 11, True, UNIST)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tline = [(title, 12, False, FG)]
    if tag:
        tline.append(("   " + tag, 10.5, True, HDC))
    text(s, Inches(2.95), y + Inches(0.07), Inches(9.6), Inches(0.5), [tline],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y += Inches(0.66)

# ===========================================================================
# 11. PATENTS
# ===========================================================================
s = slide()
chrome(s, "PATENTS", "특허", BRAND, sub="등록 6건 · 출원 1건")
pats = [
    ("KR 10-2809119", "인공지능 모델 기반 실내 건축 공사 하자 탐지 장치"),
    ("KR 10-2857163", "인공신경망 복합모델을 이용한 주차장 화재 탐지 방법·시스템"),
    ("KR 10-2876730", "사용자로그를 고려한 엘리베이터 내부 반려동물 유무 탐지"),
    ("KR 10-2910653", "인공신경망을 이용한 건설현장 충돌위험 탐지 방법·시스템"),
    ("KR 10-2958364", "인공신경망을 이용한 공간 점유율 추정 방법·시스템"),
    ("KR 10-2971047", "행동인식 기반 IoT 가전기기 제어 방법·시스템"),
    ("출원 중", "적합성 높은 스타일변환이미지 생성 방법·시스템"),
]
y = Inches(2.15)
for no, title in pats:
    reg = no != "출원 중"
    col = HDC if reg else SOFTGRAY
    rect(s, Inches(0.7), y + Inches(0.03), Inches(2.4), Inches(0.44),
         fill=BG_SOFT, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.7), y + Inches(0.03), Inches(2.4), Inches(0.44),
         [[(no, 10.5, True, col)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.35), y + Inches(0.05), Inches(9.2), Inches(0.45),
         [[(title, 12, False, FG)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.64)

# ===========================================================================
# 12. AWARDS & CERTS
# ===========================================================================
s = slide()
chrome(s, "AWARDS & CERTIFICATIONS", "수상 · 어학 · 자격", BRAND)
cols = [
    ("수상", HDC, ["CVPR'22 LOVEU AQTC Challenge — 2nd Place (CVPRW 2022)",
                  "교내 영어토론대회 3등 우수상 (2017.11)"]),
    ("어학", UNIST, ["TOEIC 965", "TOEIC Speaking 170 (Advanced Low)"]),
    ("자격증", NCC, ["정보처리기사 (2022.06.02)", "자격번호 21201071305E"]),
]
x = Inches(0.7); cw = Inches(3.85); gap = Inches(0.18)
for i, (head, col, items) in enumerate(cols):
    xx = x + i * (cw + gap)
    card(s, xx, Inches(2.3), cw, Inches(3.6))
    rect(s, xx, Inches(2.3), cw, Inches(0.12), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, xx + Inches(0.3), Inches(2.6), cw - Inches(0.6), Inches(0.5),
         [[(head, 17, True, col)]])
    yy = Inches(3.25)
    for it in items:
        text(s, xx + Inches(0.3), yy, cw - Inches(0.5), Inches(0.7),
             [[(it, 12, False, MUTED)]], line_spacing=1.18)
        yy += Inches(0.82)

# ===========================================================================
# 13. CLOSING
# ===========================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=BRAND)
text(s, Inches(0), Inches(2.6), EMU_W, Inches(1.0),
     [[("감사합니다", 40, True, WHITE)]], align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(3.7), EMU_W, Inches(0.5),
     [[("김종은 Jongeun Kim · AI Engineer", 16, False, (210, 214, 220))]],
     align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(4.4), EMU_W, Inches(0.8),
     [[("github.com/joannekim0420", 13, False, (180, 186, 196))],
      [("joannekim0420.github.io   ·   JongeunKim@r114.com", 13, False, (180, 186, 196))]],
     align=PP_ALIGN.CENTER, line_spacing=1.4, space_after=4)

out = os.path.join(ROOT, "Jongeun_Kim_Portfolio.pptx")
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
